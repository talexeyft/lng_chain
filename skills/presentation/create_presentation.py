import argparse
import json
from pathlib import Path

from pptx import Presentation


def build_presentation(slides: list[dict], output_path: Path) -> Path:
    presentation = Presentation()

    for index, slide_data in enumerate(slides):
        title = slide_data.get("title", f"Слайд {index + 1}")
        bullets = slide_data.get("bullets", [])

        if index == 0:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = title
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get("subtitle", "Автоматически создано Deep Agent")
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()

        for bullet_index, bullet in enumerate(bullets):
            if bullet_index == 0:
                text_frame.paragraphs[0].text = str(bullet)
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = str(bullet)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return output_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Генерация .pptx презентации из JSON.")
    parser.add_argument("--input", required=True, help="Путь к JSON со слайдами.")
    parser.add_argument("--output", required=True, help="Путь для выходного .pptx файла.")
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output)

    payload = json.loads(input_path.read_text(encoding="utf-8"))
    slides = payload.get("slides", [])

    if not slides:
        raise ValueError("Поле slides не должно быть пустым.")

    result_path = build_presentation(slides=slides, output_path=output_path)
    print(f"Presentation created: {result_path}")


if __name__ == "__main__":
    main()
